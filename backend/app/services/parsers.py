"""Document parsing, cleaning, and deterministic text splitting."""
from __future__ import annotations

import csv
import io
import re
import unicodedata
from pathlib import Path
from typing import Any


SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".docx", ".csv", ".xlsx", ".pptx"}

# Prefer strict UTF-8, then common Chinese Windows encodings (GB18030 ⊃ GBK/GB2312).
_TEXT_ENCODINGS = (
    "utf-8-sig",
    "utf-8",
    "utf-16",
    "utf-16-le",
    "utf-16-be",
    "gb18030",
    "gbk",
    "cp936",
    "big5",
)


def decode_text_bytes(raw: bytes) -> str:
    """Decode text file bytes with encoding auto-detection (fixes GBK TXT mojibake)."""
    if not raw:
        return ""

    if raw.startswith(b"\xef\xbb\xbf"):
        return raw.decode("utf-8-sig")
    if raw.startswith((b"\xff\xfe", b"\xfe\xff")):
        return raw.decode("utf-16")

    best_text = ""
    best_score = float("-inf")
    for encoding in _TEXT_ENCODINGS:
        try:
            text = raw.decode(encoding)
        except UnicodeDecodeError:
            continue
        score = _encoding_score(text)
        if score > best_score:
            best_score = score
            best_text = text
    if best_score > float("-inf"):
        return best_text
    # Last resort: keep as much readable content as possible
    return raw.decode("utf-8", errors="replace")


def _encoding_score(text: str) -> float:
    """Higher is better: more CJK/letters, fewer replacement/control chars."""
    if not text:
        return 0.0
    replacement = text.count("\ufffd")
    control = sum(1 for ch in text if ord(ch) < 32 and ch not in "\t\n\r")
    cjk = sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff")
    printable = sum(1 for ch in text if ch.isprintable() or ch in "\t\n\r")
    # Penalize classic UTF-8-as-Latin1 mojibake markers when another decode is valid
    mojibake = text.count("Ã") + text.count("Â") + text.count("å") + text.count("æ")
    return printable + cjk * 3 - replacement * 80 - control * 20 - mojibake * 2


def read_text_file(path: Path) -> str:
    return decode_text_bytes(path.read_bytes())


def parse_document(path: Path) -> list[tuple[str, dict[str, Any]]]:
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"不支持的文件类型：{ext or '未知'}")
    if ext in {".txt", ".md"}:
        return [(read_text_file(path), {})]
    if ext == ".csv":
        text = read_text_file(path)
        with io.StringIO(text, newline="") as handle:
            rows = ["\t".join(row) for row in csv.reader(handle)]
        return [("\n".join(rows), {"sheet": path.stem})]
    if ext == ".pdf":
        from pypdf import PdfReader

        return [
            (page.extract_text() or "", {"page": number})
            for number, page in enumerate(PdfReader(str(path)).pages, start=1)
        ]
    if ext == ".docx":
        from docx import Document

        doc = Document(str(path))
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        table_text = []
        for table in doc.tables:
            table_text.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
        return [("\n".join([text, *table_text]), {})]
    if ext == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        sections = []
        for sheet in workbook.worksheets:
            rows = [
                "\t".join("" if value is None else str(value) for value in row)
                for row in sheet.iter_rows(values_only=True)
            ]
            sections.append(("\n".join(rows), {"sheet": sheet.title}))
        workbook.close()
        return sections
    if ext == ".pptx":
        from pptx import Presentation

        sections = []
        for number, slide in enumerate(Presentation(str(path)).slides, start=1):
            text = "\n".join(
                shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)
            )
            sections.append((text, {"page": number}))
        return sections
    raise AssertionError("unreachable")

def clean_text(text: str) -> str:
    text = unicodedata.normalize("NFKC", text or "")
    text = text.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+", " ", line).strip() for line in text.split("\n")]
    return re.sub(r"\n{3,}", "\n\n", "\n".join(lines)).strip()


def split_text(
    text: str,
    *,
    mode: str = "recursive",
    chunk_size: int = 500,
    overlap: int = 50,
    min_chars: int = 50,
) -> list[str]:
    text = clean_text(text)
    if not text:
        return []
    chunk_size = max(1, chunk_size)
    overlap = max(0, min(overlap, chunk_size - 1))
    mode = mode if mode in {"fixed", "paragraph", "recursive"} else "recursive"

    if mode == "fixed":
        chunks = _fixed_split(text, chunk_size, overlap)
    elif mode == "paragraph":
        chunks = _join_parts(re.split(r"\n\s*\n", text), chunk_size, overlap)
    else:
        parts = re.split(r"(?<=[。！？.!?])\s*|\n{2,}|\n", text)
        chunks = _join_parts(parts, chunk_size, overlap)

    filtered = [chunk.strip() for chunk in chunks if chunk.strip()]
    if len(filtered) > 1 and len(filtered[-1]) < min_chars:
        filtered[-2] = (filtered[-2] + "\n" + filtered[-1]).strip()
        filtered.pop()
    return filtered


def _fixed_split(text: str, size: int, overlap: int) -> list[str]:
    step = max(1, size - overlap)
    return [text[start : start + size] for start in range(0, len(text), step)]


def _join_parts(parts: list[str], size: int, overlap: int) -> list[str]:
    chunks: list[str] = []
    current = ""
    for raw_part in parts:
        part = raw_part.strip()
        if not part:
            continue
        if len(part) > size:
            if current:
                chunks.append(current)
                current = ""
            chunks.extend(_fixed_split(part, size, overlap))
            continue
        candidate = f"{current}\n{part}".strip() if current else part
        if len(candidate) <= size:
            current = candidate
            continue
        chunks.append(current)
        prefix = current[-overlap:] if overlap else ""
        current = f"{prefix}\n{part}".strip()
        if len(current) > size:
            chunks.extend(_fixed_split(current, size, overlap)[:-1])
            current = _fixed_split(current, size, overlap)[-1]
    if current:
        chunks.append(current)
    return chunks
