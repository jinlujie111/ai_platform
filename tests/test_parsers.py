from backend.app.services.parsers import clean_text, parse_document, split_text


def test_clean_text_normalizes_whitespace():
    assert clean_text("Ａ  B\r\n\r\n\r\n C\x00") == "A B\n\nC"


def test_fixed_split_respects_overlap_and_merges_short_tail():
    chunks = split_text(
        "abcdefghijklmnopqrstuvwxyz",
        mode="fixed",
        chunk_size=10,
        overlap=2,
        min_chars=5,
    )
    assert chunks[0] == "abcdefghij"
    assert chunks[1].startswith("ijkl")
    assert chunks[-1].endswith("z")


def test_recursive_split_keeps_content():
    text = "第一句。第二句！\n\n第三段内容。"
    chunks = split_text(text, mode="recursive", chunk_size=12, overlap=2, min_chars=1)
    assert "第一句。" in "".join(chunks)
    assert "第三段内容。" in "".join(chunks)


def test_supported_document_parsers(tmp_path):
    txt = tmp_path / "sample.txt"
    txt.write_text("纯文本", encoding="utf-8")
    csv_file = tmp_path / "sample.csv"
    csv_file.write_text("name,value\n测试,1", encoding="utf-8")

    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from pypdf import PdfWriter

    docx_file = tmp_path / "sample.docx"
    doc = Document()
    doc.add_paragraph("Word 内容")
    doc.save(docx_file)

    xlsx_file = tmp_path / "sample.xlsx"
    workbook = Workbook()
    workbook.active.append(["表格", 2])
    workbook.save(xlsx_file)

    pptx_file = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "幻灯片内容"
    presentation.save(pptx_file)

    pdf_file = tmp_path / "sample.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with pdf_file.open("wb") as handle:
        writer.write(handle)

    assert parse_document(txt)[0][0] == "纯文本"
    assert "测试" in parse_document(csv_file)[0][0]
    assert "Word 内容" in parse_document(docx_file)[0][0]
    assert "表格" in parse_document(xlsx_file)[0][0]
    assert "幻灯片内容" in parse_document(pptx_file)[0][0]
    assert parse_document(pdf_file)[0][1]["page"] == 1


def test_txt_decodes_gbk_without_mojibake(tmp_path):
    from backend.app.services.parsers import decode_text_bytes

    sample = "知识库配置上传中文内容"
    assert decode_text_bytes(sample.encode("gbk")) == sample
    assert decode_text_bytes(sample.encode("gb18030")) == sample
    assert decode_text_bytes(sample.encode("utf-8")) == sample

    gbk_file = tmp_path / "gbk.txt"
    gbk_file.write_bytes(sample.encode("gbk"))
    assert parse_document(gbk_file)[0][0] == sample

    csv_gbk = tmp_path / "gbk.csv"
    csv_gbk.write_bytes("name,value\n测试,1".encode("gbk"))
    assert "测试" in parse_document(csv_gbk)[0][0]
